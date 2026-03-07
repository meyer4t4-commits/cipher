#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Color palette
DARK_BG = RGBColor(10, 14, 26)          # Very dark navy
CARD_BG = RGBColor(27, 58, 92)          # Deep navy
ACCENT_BLUE = RGBColor(76, 110, 245)    # Electric blue
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(224, 224, 224)
MUTED_GRAY = RGBColor(160, 160, 160)
ACCENT_GREEN = RGBColor(16, 185, 129)   # Green accent

def add_slide_background(slide, color=DARK_BG):
    """Add solid background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide():
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_slide_background(slide)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "FORGE"
    title_p.font.size = Pt(72)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(0.5))
    tagline_frame = tagline_box.text_frame
    tagline_p = tagline_frame.paragraphs[0]
    tagline_p.text = "Your Business. Your AI. Your Advantage."
    tagline_p.font.size = Pt(28)
    tagline_p.font.color.rgb = ACCENT_BLUE
    tagline_p.alignment = PP_ALIGN.CENTER

    # Subheading
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.4))
    sub_frame = sub_box.text_frame
    sub_p = sub_frame.paragraphs[0]
    sub_p.text = "Powered by Elysian Protocol"
    sub_p.font.size = Pt(16)
    sub_p.font.color.rgb = MUTED_GRAY
    sub_p.alignment = PP_ALIGN.CENTER

    # Accent line
    line = slide.shapes.add_shape(1, Inches(3.5), Inches(4.2), Inches(3), Inches(0))
    line.line.color.rgb = ACCENT_BLUE
    line.line.width = Pt(2)
    line.fill.background()

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.35))
    contact_frame = contact_box.text_frame
    contact_p = contact_frame.paragraphs[0]
    contact_p.text = "Mark Meyer  |  elysianprotocol.io"
    contact_p.font.size = Pt(12)
    contact_p.font.color.rgb = MUTED_GRAY
    contact_p.alignment = PP_ALIGN.CENTER

def add_problem_slide():
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "The Problem"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    problems = [
        ("Generic AI Tools", "ChatGPT doesn't know your business, your customers, or your workflows"),
        ("Data Privacy Concerns", "Your business information fed to public AI providers"),
        ("Manual, Repetitive Work", "Teams spend hours on tasks that should be automated"),
        ("No Memory", "AI forgets context between sessions—you have to repeat yourself")
    ]

    y_pos = 1.3
    for title, desc in problems:
        # Background card
        card = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(4.2), Inches(0.95))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG

        # Accent left border
        accent = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(0.08), Inches(0.95))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT_BLUE
        accent.line.color.rgb = ACCENT_BLUE

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.08), Inches(3.8), Inches(0.3))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.color.rgb = ACCENT_BLUE

        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.42), Inches(3.8), Inches(0.45))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_p = desc_frame.paragraphs[0]
        desc_p.text = desc
        desc_p.font.size = Pt(11)
        desc_p.font.color.rgb = LIGHT_GRAY

        y_pos += 1.05

def add_solution_slide():
    """Slide 3: The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "The Solution: Forge"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    # Left side description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.8), Inches(3.2))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True

    bullets = [
        ("Your own AI intelligence", True),
        ("Customized for your business operations", False),
        ("Remembers every interaction and decision", False),
        ("Routes to the best model for each task", False),
        ("Your data stays completely private", False),
        ("Learns your business rules and workflows", False)
    ]

    for idx, (text, is_bold) in enumerate(bullets):
        if idx == 0:
            p = desc_frame.paragraphs[0]
        else:
            p = desc_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = is_bold
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(0)
        p.space_after = Pt(6)

    # Right side card
    card = slide.shapes.add_shape(1, Inches(5.7), Inches(1.3), Inches(3.8), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BG

    card_title = slide.shapes.add_textbox(Inches(5.9), Inches(1.45), Inches(3.4), Inches(0.3))
    card_title_frame = card_title.text_frame
    card_title_p = card_title_frame.paragraphs[0]
    card_title_p.text = "Sovereign AI"
    card_title_p.font.size = Pt(18)
    card_title_p.font.bold = True
    card_title_p.font.color.rgb = ACCENT_BLUE

    card_text = slide.shapes.add_textbox(Inches(5.9), Inches(1.85), Inches(3.4), Inches(0.8))
    card_text_frame = card_text.text_frame
    card_text_frame.word_wrap = True
    card_text_p = card_text_frame.paragraphs[0]
    card_text_p.text = "Your business intelligence stays in your control. Persistent memory. Multi-model routing. Zero data sharing."
    card_text_p.font.size = Pt(12)
    card_text_p.font.color.rgb = LIGHT_GRAY

    # Stats section
    stats_box = slide.shapes.add_shape(1, Inches(5.7), Inches(3.0), Inches(3.8), Inches(1.5))
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = ACCENT_BLUE
    stats_box.line.color.rgb = ACCENT_BLUE

    stat_num = slide.shapes.add_textbox(Inches(5.9), Inches(3.15), Inches(1.7), Inches(0.5))
    stat_num_frame = stat_num.text_frame
    stat_num_p = stat_num_frame.paragraphs[0]
    stat_num_p.text = "100%"
    stat_num_p.font.size = Pt(48)
    stat_num_p.font.bold = True
    stat_num_p.font.color.rgb = DARK_BG
    stat_num_p.alignment = PP_ALIGN.CENTER

    stat_label = slide.shapes.add_textbox(Inches(5.9), Inches(3.75), Inches(1.7), Inches(0.6))
    stat_label_frame = stat_label.text_frame
    stat_label_frame.word_wrap = True
    stat_label_p = stat_label_frame.paragraphs[0]
    stat_label_p.text = "Data Privacy"
    stat_label_p.font.size = Pt(13)
    stat_label_p.font.bold = True
    stat_label_p.font.color.rgb = DARK_BG
    stat_label_p.alignment = PP_ALIGN.CENTER

    stat_desc = slide.shapes.add_textbox(Inches(7.7), Inches(3.15), Inches(1.65), Inches(1.2))
    stat_desc_frame = stat_desc.text_frame
    stat_desc_frame.word_wrap = True
    stat_desc_p = stat_desc_frame.paragraphs[0]
    stat_desc_p.text = "Your data. Never shared."
    stat_desc_p.font.size = Pt(12)
    stat_desc_p.font.color.rgb = DARK_BG
    stat_desc_p.alignment = PP_ALIGN.CENTER

def add_how_it_works_slide():
    """Slide 4: How It Works"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "How It Works"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    steps = [
        ("1", "Connect", "Link your business data, tools, and workflows"),
        ("2", "Configure", "Define rules, playbooks, and AI behaviors"),
        ("3", "Command", "Deploy Forge to automate and amplify your team")
    ]

    x_positions = [0.8, 3.6, 6.4]

    for idx, (num, step_title, desc) in enumerate(steps):
        x = x_positions[idx]

        # Circle number
        circle = slide.shapes.add_shape(3, Inches(x), Inches(1.3), Inches(0.7), Inches(0.7))  # 3 = oval
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.color.rgb = ACCENT_BLUE

        num_box = slide.shapes.add_textbox(Inches(x), Inches(1.3), Inches(0.7), Inches(0.7))
        num_frame = num_box.text_frame
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        num_p = num_frame.paragraphs[0]
        num_p.text = num
        num_p.font.size = Pt(36)
        num_p.font.bold = True
        num_p.font.color.rgb = DARK_BG
        num_p.alignment = PP_ALIGN.CENTER

        # Title
        step_title_box = slide.shapes.add_textbox(Inches(x - 0.15), Inches(2.15), Inches(1.0), Inches(0.4))
        step_title_frame = step_title_box.text_frame
        step_title_p = step_title_frame.paragraphs[0]
        step_title_p.text = step_title
        step_title_p.font.size = Pt(18)
        step_title_p.font.bold = True
        step_title_p.font.color.rgb = WHITE
        step_title_p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x - 0.3), Inches(2.65), Inches(1.6), Inches(1.2))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_p = desc_frame.paragraphs[0]
        desc_p.text = desc
        desc_p.font.size = Pt(12)
        desc_p.font.color.rgb = LIGHT_GRAY
        desc_p.alignment = PP_ALIGN.CENTER

    # Bottom CTA
    cta_box = slide.shapes.add_shape(1, Inches(1.5), Inches(4.3), Inches(7), Inches(1.0))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = CARD_BG
    cta_box.line.color.rgb = CARD_BG

    cta_text = slide.shapes.add_textbox(Inches(1.7), Inches(4.4), Inches(6.6), Inches(0.8))
    cta_frame = cta_text.text_frame
    cta_frame.word_wrap = True
    cta_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cta_p = cta_frame.paragraphs[0]
    cta_p.text = "It's that simple. Forge becomes an extension of your team in minutes."
    cta_p.font.size = Pt(14)
    cta_p.font.color.rgb = ACCENT_BLUE
    cta_p.alignment = PP_ALIGN.CENTER

def add_technology_slide():
    """Slide 5: The Technology"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "The Technology: Orchid"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.35))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "Multi-Model AI Router Built for Performance"
    subtitle_p.font.size = Pt(16)
    subtitle_p.font.color.rgb = ACCENT_BLUE

    # Models
    models = [
        ("Claude 3.5", "Strategic thinking, reasoning, complex decisions", ACCENT_BLUE),
        ("Llama 3", "Speed and efficiency, real-time responses", ACCENT_GREEN),
        ("DeepSeek", "Code generation, technical problem solving", ACCENT_BLUE)
    ]

    x_positions = [0.5, 3.65, 6.8]
    for idx, (model_name, specialty, accent_color) in enumerate(models):
        x = x_positions[idx]

        # Card
        card = slide.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(2.9), Inches(3.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG

        # Top accent
        top_accent = slide.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(2.9), Inches(0.08))
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = accent_color
        top_accent.line.color.rgb = accent_color

        # Name
        name_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.0), Inches(2.5), Inches(0.4))
        name_frame = name_box.text_frame
        name_p = name_frame.paragraphs[0]
        name_p.text = model_name
        name_p.font.size = Pt(18)
        name_p.font.bold = True
        name_p.font.color.rgb = accent_color

        # Specialty
        spec_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.6), Inches(2.5), Inches(2.0))
        spec_frame = spec_box.text_frame
        spec_frame.word_wrap = True
        spec_p = spec_frame.paragraphs[0]
        spec_p.text = specialty
        spec_p.font.size = Pt(12)
        spec_p.font.color.rgb = LIGHT_GRAY

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.word_wrap = True
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "Intelligent routing ensures every task reaches the perfect model—no wasted compute, no suboptimal results."
    footer_p.font.size = Pt(13)
    footer_p.font.color.rgb = MUTED_GRAY

def add_differentiators_slide():
    """Slide 6: Why Forge?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Why Forge? The Cipher Advantage"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    differentiators = [
        ("⚡", "Persistent Memory", "Remembers every interaction, decision, and preference across all sessions"),
        ("🔐", "Sovereign Data", "Zero external data sharing. Your business intelligence stays completely private"),
        ("🧠", "Business Customization", "Learns your workflows, rules, and company playbooks automatically"),
        ("⚙️", "Multi-Model Intelligence", "Routes to Claude, Llama, or DeepSeek based on task requirements"),
        ("🎯", "Purpose-Built for SMBs", "Designed for roofing, construction, restaurants, accounting—vertical-specific"),
        ("💰", "Premium Simplicity", "$99-$499/mo. More intelligent. More private. More yours.")
    ]

    y_pos = 1.2
    for idx, (emoji, title, desc) in enumerate(differentiators):
        is_left = idx % 2 == 0
        x_pos = 0.5 if is_left else 5.2

        if idx % 2 == 0 and idx > 0:
            y_pos += 1.0

        # Card
        card = slide.shapes.add_shape(1, Inches(x_pos), Inches(y_pos), Inches(4.3), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG

        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(x_pos + 0.15), Inches(y_pos + 0.15), Inches(0.5), Inches(0.6))
        emoji_frame = emoji_box.text_frame
        emoji_p = emoji_frame.paragraphs[0]
        emoji_p.text = emoji
        emoji_p.font.size = Pt(28)

        # Title
        title_box = slide.shapes.add_textbox(Inches(x_pos + 0.75), Inches(y_pos + 0.1), Inches(3.5), Inches(0.3))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(13)
        title_p.font.bold = True
        title_p.font.color.rgb = ACCENT_BLUE

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x_pos + 0.75), Inches(y_pos + 0.45), Inches(3.5), Inches(0.35))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_p = desc_frame.paragraphs[0]
        desc_p.text = desc
        desc_p.font.size = Pt(10)
        desc_p.font.color.rgb = LIGHT_GRAY

def add_use_cases_slide():
    """Slide 7: Use Cases"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Built for Your Business"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    use_cases = [
        ("Roofing & Construction", [
            "Estimate generation from photos & specs",
            "Customer follow-up automation",
            "Project scheduling and crew optimization"
        ]),
        ("Restaurants & Hospitality", [
            "Inventory forecasting and ordering",
            "Staff scheduling based on demand",
            "Customer feedback analysis & reviews"
        ]),
        ("Accounting & Finance", [
            "Invoice processing and categorization",
            "Tax compliance research automation",
            "Financial report generation"
        ])
    ]

    y_pos = 1.3
    for industry, examples in use_cases:
        # Background card
        card = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG

        # Accent left border
        accent = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(0.08), Inches(1.1))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT_BLUE
        accent.line.color.rgb = ACCENT_BLUE

        # Industry title
        industry_box = slide.shapes.add_textbox(Inches(0.75), Inches(y_pos + 0.08), Inches(8.7), Inches(0.3))
        industry_frame = industry_box.text_frame
        industry_p = industry_frame.paragraphs[0]
        industry_p.text = industry
        industry_p.font.size = Pt(15)
        industry_p.font.bold = True
        industry_p.font.color.rgb = ACCENT_BLUE

        # Examples
        examples_box = slide.shapes.add_textbox(Inches(0.95), Inches(y_pos + 0.42), Inches(8.5), Inches(0.6))
        examples_frame = examples_box.text_frame
        examples_frame.word_wrap = True

        for i, example in enumerate(examples):
            if i == 0:
                p = examples_frame.paragraphs[0]
            else:
                p = examples_frame.add_paragraph()
            p.text = example
            p.level = 0
            p.font.size = Pt(11)
            p.font.color.rgb = LIGHT_GRAY
            p.space_before = Pt(0)
            p.space_after = Pt(2)

        y_pos += 1.25

def add_pilot_slide():
    """Slide 8: Pilot Program"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Join the Pilot"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    # Main CTA
    cta_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.3), Inches(7), Inches(1.5))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = ACCENT_BLUE
    cta_box.line.color.rgb = ACCENT_BLUE

    cta_main = slide.shapes.add_textbox(Inches(1.7), Inches(1.5), Inches(6.6), Inches(0.4))
    cta_main_frame = cta_main.text_frame
    cta_main_p = cta_main_frame.paragraphs[0]
    cta_main_p.text = "30-Day Free Pilot"
    cta_main_p.font.size = Pt(36)
    cta_main_p.font.bold = True
    cta_main_p.font.color.rgb = DARK_BG
    cta_main_p.alignment = PP_ALIGN.CENTER

    cta_desc = slide.shapes.add_textbox(Inches(1.7), Inches(2.0), Inches(6.6), Inches(0.7))
    cta_desc_frame = cta_desc.text_frame
    cta_desc_frame.word_wrap = True
    cta_desc_p = cta_desc_frame.paragraphs[0]
    cta_desc_p.text = "Full access to Forge and the Orchid engine. For early partners ready to build the future of AI-powered business."
    cta_desc_p.font.size = Pt(14)
    cta_desc_p.font.color.rgb = DARK_BG
    cta_desc_p.alignment = PP_ALIGN.CENTER

    # What's included
    included_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.3))
    included_title_frame = included_title.text_frame
    included_title_p = included_title_frame.paragraphs[0]
    included_title_p.text = "What's Included:"
    included_title_p.font.size = Pt(16)
    included_title_p.font.bold = True
    included_title_p.font.color.rgb = ACCENT_BLUE

    # Pilot items
    pilot_items = [
        "Unlimited Forge deployments for your team",
        "Direct support from the Elysian Protocol team",
        "Custom integration with your existing tools",
        "Onboarding and training for your workflows"
    ]

    items_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1.6))
    items_frame = items_box.text_frame
    items_frame.word_wrap = True

    for i, item in enumerate(pilot_items):
        if i == 0:
            p = items_frame.paragraphs[0]
        else:
            p = items_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(0)
        p.space_after = Pt(4)

def add_traction_slide():
    """Slide 9: Traction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Traction"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    traction_items = [
        ("Founder-Led", "Solo founder building with deep technical expertise"),
        ("Orchid Engine Live", "Multi-model routing fully operational and tested"),
        ("Integration Ready", "Connectors for Zapier, Make, custom APIs"),
        ("Customer Validation", "Early conversations with roofing, construction, restaurant owners")
    ]

    y_pos = 1.4
    for metric, detail in traction_items:
        # Metric box
        metric_box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(4.3), Inches(0.85))
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = CARD_BG
        metric_box.line.color.rgb = CARD_BG

        metric_title = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.08), Inches(3.9), Inches(0.3))
        metric_title_frame = metric_title.text_frame
        metric_title_p = metric_title_frame.paragraphs[0]
        metric_title_p.text = metric
        metric_title_p.font.size = Pt(15)
        metric_title_p.font.bold = True
        metric_title_p.font.color.rgb = ACCENT_BLUE

        metric_desc = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.43), Inches(3.9), Inches(0.35))
        metric_desc_frame = metric_desc.text_frame
        metric_desc_frame.word_wrap = True
        metric_desc_p = metric_desc_frame.paragraphs[0]
        metric_desc_p.text = detail
        metric_desc_p.font.size = Pt(11)
        metric_desc_p.font.color.rgb = LIGHT_GRAY

        # Accent dot
        dot = slide.shapes.add_shape(3, Inches(9.0), Inches(y_pos + 0.32), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_GREEN
        dot.line.color.rgb = ACCENT_GREEN

        y_pos += 1.0

    # Vision box
    vision_box = slide.shapes.add_shape(1, Inches(5.2), Inches(1.4), Inches(4.3), Inches(3.4))
    vision_box.fill.solid()
    vision_box.fill.fore_color.rgb = CARD_BG
    vision_box.line.color.rgb = CARD_BG

    vision_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(3.9), Inches(0.3))
    vision_title_frame = vision_title.text_frame
    vision_title_p = vision_title_frame.paragraphs[0]
    vision_title_p.text = "The Vision"
    vision_title_p.font.size = Pt(16)
    vision_title_p.font.bold = True
    vision_title_p.font.color.rgb = ACCENT_BLUE

    vision_text = slide.shapes.add_textbox(Inches(5.4), Inches(2.0), Inches(3.9), Inches(2.6))
    vision_text_frame = vision_text.text_frame
    vision_text_frame.word_wrap = True
    vision_text_p = vision_text_frame.paragraphs[0]
    vision_text_p.text = "Every business should have its own AI intelligence. No giant corporations should own your business logic. No vendor lock-in. No data selling. Sovereign AI for the modern business."
    vision_text_p.font.size = Pt(12)
    vision_text_p.font.color.rgb = LIGHT_GRAY

def add_contact_slide():
    """Slide 10: Contact & CTA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    # Large accent background on right
    accent_bg = slide.shapes.add_shape(1, Inches(5.5), Inches(0), Inches(4.5), Inches(5.625))
    accent_bg.fill.solid()
    accent_bg.fill.fore_color.rgb = CARD_BG
    accent_bg.line.color.rgb = CARD_BG

    # Left side - main text
    main_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(0.6))
    main_title_frame = main_title.text_frame
    main_title_frame.word_wrap = True
    main_title_p = main_title_frame.paragraphs[0]
    main_title_p.text = "Let's Build the Future"
    main_title_p.font.size = Pt(44)
    main_title_p.font.bold = True
    main_title_p.font.color.rgb = WHITE

    main_subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5), Inches(0.4))
    main_subtitle_frame = main_subtitle.text_frame
    main_subtitle_p = main_subtitle_frame.paragraphs[0]
    main_subtitle_p.text = "Sovereign AI for your business starts today."
    main_subtitle_p.font.size = Pt(18)
    main_subtitle_p.font.color.rgb = ACCENT_BLUE

    # Right side - contact info
    contact_name = slide.shapes.add_textbox(Inches(5.7), Inches(1.3), Inches(4.1), Inches(0.4))
    contact_name_frame = contact_name.text_frame
    contact_name_p = contact_name_frame.paragraphs[0]
    contact_name_p.text = "Mark Meyer"
    contact_name_p.font.size = Pt(28)
    contact_name_p.font.bold = True
    contact_name_p.font.color.rgb = WHITE

    contact_title = slide.shapes.add_textbox(Inches(5.7), Inches(1.8), Inches(4.1), Inches(0.3))
    contact_title_frame = contact_title.text_frame
    contact_title_p = contact_title_frame.paragraphs[0]
    contact_title_p.text = "Founder, Elysian Protocol"
    contact_title_p.font.size = Pt(13)
    contact_title_p.font.color.rgb = ACCENT_BLUE

    # Contact details
    contact_details = [
        ("Email", "mark@elysianprotocol.io"),
        ("Web", "elysianprotocol.io"),
        ("Twitter", "@markmeyeragi")
    ]

    contact_y = 2.4
    for label, value in contact_details:
        label_box = slide.shapes.add_textbox(Inches(5.7), Inches(contact_y), Inches(4.1), Inches(0.2))
        label_frame = label_box.text_frame
        label_p = label_frame.paragraphs[0]
        label_p.text = label
        label_p.font.size = Pt(10)
        label_p.font.color.rgb = MUTED_GRAY

        value_box = slide.shapes.add_textbox(Inches(5.7), Inches(contact_y + 0.25), Inches(4.1), Inches(0.3))
        value_frame = value_box.text_frame
        value_p = value_frame.paragraphs[0]
        value_p.text = value
        value_p.font.size = Pt(12)
        value_p.font.bold = True
        value_p.font.color.rgb = WHITE

        contact_y += 0.7

    # CTA button
    cta_btn = slide.shapes.add_shape(1, Inches(5.7), Inches(4.0), Inches(4.1), Inches(0.5))
    cta_btn.fill.solid()
    cta_btn.fill.fore_color.rgb = ACCENT_BLUE
    cta_btn.line.color.rgb = ACCENT_BLUE

    cta_btn_text = slide.shapes.add_textbox(Inches(5.7), Inches(4.0), Inches(4.1), Inches(0.5))
    cta_btn_frame = cta_btn_text.text_frame
    cta_btn_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cta_btn_p = cta_btn_frame.paragraphs[0]
    cta_btn_p.text = "Start Your Pilot Today"
    cta_btn_p.font.size = Pt(14)
    cta_btn_p.font.bold = True
    cta_btn_p.font.color.rgb = DARK_BG
    cta_btn_p.alignment = PP_ALIGN.CENTER

# Build the presentation
add_title_slide()
add_problem_slide()
add_solution_slide()
add_how_it_works_slide()
add_technology_slide()
add_differentiators_slide()
add_use_cases_slide()
add_pilot_slide()
add_traction_slide()
add_contact_slide()

# Save
output_path = "/sessions/inspiring-funny-rubin/mnt/orchid/FORGE_PITCH_DECK.pptx"
prs.save(output_path)
print(f"Forge pitch deck created successfully at {output_path}")
